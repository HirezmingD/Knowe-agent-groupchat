# v1.0.19.4 — 用户上传附件：凭证护栏 + 多模态打包 + 注入
"""
attachments.py — 把「用户选/拖进来的本地文件」安全地喂给 LLM。

这里只做三件小事，一件都不多（DESIGN §七「用简单的代码」）：

  1. **凭证护栏（DESIGN 决策 #9）**。后端直读本地路径，天生有「被诱导读任意文件」
     的风险：user_message 里的 path 只是一个字符串，谁都能塞 C:\\Windows\\... 或 /etc/passwd。
     护栏是一枚 HMAC 签名：只有 Electron 主进程（它亲眼看着用户点了对话框 / 把文件拖进
     窗口）用**主进程与后端共享的 runtime_token** 对绝对路径签名，后端才认。
     消息正文里凭空捏造的路径没有合法签名 → 一律拒读。签的是**路径本身**，
     所以同一份用户选过的文件，回看时重新读取也天然放行（DESIGN 决策 #3）。

  2. **原样打包（DESIGN 决策 #4/#5）**。图片 → OpenAI 多模态 image_url（data URL）；
     其余 → OpenAI file 内容块（base64 file_data）。**只做 base64 装信封，不转码、不 OCR、
     不缩图**——字节原样。模型收不收由模型自己决定（收不到就 HTTP 400，上层友好打回）。

  3. **注入（本次故障根因的修复点）**。用户消息在后端全程是一个**纯字符串 content**，
     两个 Agent 都把它拼成 {"role":"user","content": <str>} 发给 provider——附件在这条链路上
     根本没有落脚点，于是「界面层完整、LLM 输入层丢失」。这里提供把「文本 + 附件块」
     合成 OpenAI 多模态数组的工具，在**投影之后、发 provider 之前**替换掉最后一条 user 消息。
     只影响这一回合发出去的投影副本，Agent 的权威历史仍存纯文本（不把 base64 塞进历史/落盘）。
"""

from __future__ import annotations

import base64
import hashlib
import hmac
import io
import logging
import mimetypes
import os
from typing import Any

log = logging.getLogger("knowe.attachments")

# 常见图片扩展名 → 走 image_url（视觉模型能看）。其余走 file 块。
_IMAGE_MIME_PREFIX = "image/"
_EXT_MIME_FALLBACK = {
    "md": "text/markdown",
    "markdown": "text/markdown",
    "txt": "text/plain",
    "log": "text/plain",
    "json": "application/json",
    "csv": "text/csv",
    "tsv": "text/tab-separated-values",
    "py": "text/x-python",
    "js": "text/javascript",
    "ts": "text/plain",
    "yaml": "text/plain",
    "yml": "text/plain",
    "pdf": "application/pdf",
    "webp": "image/webp",
}


class AttachmentError(Exception):
    """一个能直接说给用户听的附件错误。kind 供上层决定文案分类。"""

    def __init__(self, message: str, *, kind: str = "invalid") -> None:
        super().__init__(message)
        self.kind = kind


def normalize_path(path: Any) -> str:
    """签名与读取共用的规范绝对路径。主进程 path.resolve 与此在同机同 OS 下一致。"""
    return os.path.abspath(os.path.expanduser(str(path or "")))


def sign_path(path: Any, token: str) -> str:
    """用共享 runtime_token 对绝对路径签名（主进程与后端算法必须逐字一致）。"""
    key = (token or "").encode("utf-8")
    msg = normalize_path(path).encode("utf-8")
    return hmac.new(key, msg, hashlib.sha256).hexdigest()


def verify_path(path: Any, sig: Any, token: str) -> bool:
    """常量时间校验：这条 path 是不是主进程用 token 亲手签发的。"""
    if not token or not isinstance(sig, str) or not sig:
        return False
    try:
        expected = sign_path(path, token)
    except Exception:
        return False
    return hmac.compare_digest(expected, sig)


def _basename(path: str) -> str:
    normalized = str(path or "").replace("\\", "/").rstrip("/")
    tail = normalized.rsplit("/", 1)[-1]
    return tail or (str(path) or "file")


def _guess_mime(name: str, ext: str) -> str:
    guessed, _ = mimetypes.guess_type(name)
    if guessed:
        return guessed
    key = (ext or "").lower().lstrip(".")
    if key in _EXT_MIME_FALLBACK:
        return _EXT_MIME_FALLBACK[key]
    return "application/octet-stream"


def load_part(record: dict[str, Any], token: str) -> tuple[dict[str, Any], dict[str, Any]]:
    """把一条附件记录变成 (provider 内容块, 回显用元数据)。

    护栏在这里落地：签名不对 → 拒读（kind=guard）；文件不在 → 友好打回（kind=missing）。
    """
    if not isinstance(record, dict):
        raise AttachmentError("附件数据格式不对。", kind="invalid")
    path = record.get("path")
    name = str(record.get("name") or _basename(str(path or "")))
    ext = str(record.get("ext") or (name.rsplit(".", 1)[-1] if "." in name else ""))
    sig = record.get("sig")

    if not path or not verify_path(path, sig, token):
        # 消息正文捏造的路径走到这里——没有主进程签名，读都不读。
        raise AttachmentError(
            f"附件「{name}」没有通过安全校验，已被拒绝（只有通过「添加附件」或拖拽选进来的文件才会被发送）。",
            kind="guard",
        )
    abs_path = normalize_path(path)
    if not os.path.isfile(abs_path):
        raise AttachmentError(f"原文件「{name}」已被移动、重命名或删除。", kind="missing")

    try:
        with open(abs_path, "rb") as handle:
            data = handle.read()
    except OSError as exc:
        raise AttachmentError(f"读取附件「{name}」失败：{exc}", kind="missing") from None

    b64 = base64.b64encode(data).decode("ascii")
    mime = _guess_mime(name, ext)
    if mime.startswith(_IMAGE_MIME_PREFIX):
        part: dict[str, Any] = {
            "type": "image_url",
            "image_url": {"url": f"data:{mime};base64,{b64}"},
        }
    else:
        # OpenAI 的 file 内容块（PDF/文本等）；模型不支持就 400，上层友好打回。
        part = {
            "type": "file",
            "file": {"filename": name, "file_data": f"data:{mime};base64,{b64}"},
        }
    meta = {
        "path": abs_path,
        "name": name,
        "ext": ext,
        "size": len(data),
        "sig": sig,
    }
    return part, meta


def build_parts(
    records: list[dict[str, Any]] | None, token: str,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """读一批附件 → (provider 内容块列表, 元数据列表)。任一条失败即抛（原子性）。"""
    parts: list[dict[str, Any]] = []
    metas: list[dict[str, Any]] = []
    for record in records or []:
        part, meta = load_part(record, token)
        parts.append(part)
        metas.append(meta)
    return parts, metas


def echo_meta(record: dict[str, Any]) -> dict[str, Any]:
    """user_echo / 历史里带的附件元数据——只带路径与身份，**绝不带字节**（DESIGN 决策 #3）。"""
    path = record.get("path")
    name = str(record.get("name") or _basename(str(path or "")))
    ext = str(record.get("ext") or (name.rsplit(".", 1)[-1] if "." in name else ""))
    out: dict[str, Any] = {"path": normalize_path(path), "name": name, "ext": ext}
    size = record.get("size")
    if isinstance(size, (int, float)) and size >= 0:
        out["size"] = int(size)
    sig = record.get("sig")
    if isinstance(sig, str) and sig:
        out["sig"] = sig
    return out


def _text_of(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                return str(item.get("text") or "")
    return ""


def user_content(text: str, parts: list[dict[str, Any]] | None) -> Any:
    """无附件 → 纯文本字符串（和以前完全一样）；有附件 → OpenAI 多模态数组。"""
    if not parts:
        return text
    blocks: list[dict[str, Any]] = []
    if text:
        blocks.append({"type": "text", "text": text})
    blocks.extend(parts)
    return blocks


def inject_into_last_user(
    messages: list[dict[str, Any]], parts: list[dict[str, Any]] | None,
) -> bool:
    """把附件块并进**最后一条 user 消息**（当前回合），返回是否命中。

    只改传给 provider 的投影副本：当前回合永远是投影尾部的 verbatim 消息，
    历史里更早的 user 消息保持纯文本、不被重发 base64。
    """
    if not parts:
        return False
    for message in reversed(messages):
        if isinstance(message, dict) and message.get("role") == "user":
            message["content"] = user_content(_text_of(message.get("content")), parts)
            return True
    return False



# ═══════════════════════════════════════════════════════════════
# [v1.0.39.2] 文件块 → 文本块降级（多模态能力修复）
#
# 背景：非图片附件以 OpenAI 新版 file 块发送，兼容网关（如 opencode.ai）
# 只认 text/image_url/video_url/video，不认 file 块 → 400 整条打回。
# 这里提供三件套：
#   1. extract_file_text(part)      —— 从 file 块内 file_data(base64) 解码，
#                                      按 MIME 提取文字（fitz/docx/openpyxl/pptx）。
#   2. replace_file_blocks_with_text —— 把 messages 尾部 user 里的 file 块换成
#                                      text 块（纯函数，多附件逐个替换）。
#   3. build_format_fallback(...)    —— 给 ProviderClient 用的降级回调工厂：
#                                      命中 format_rejected="file" 时替换重发一次，
#                                      并标记该网关「不支持 file 块」能力缓存。
# ═══════════════════════════════════════════════════════════════

# 单文件提取文本上限（字符）。保护模型上下文窗口（PRD §2.2）。
EXTRACT_MAX_CHARS = 50_000

# 能力缓存：{provider|base_url|model -> "file_unsupported": True}
# 只缓存「不支持」结论（支持方零状态）；键含模型名，服务商/模型切换天然失效。
_file_unsupported: dict[str, bool] = {}


def _normalize_base(base_url: Any) -> str:
    """归一化网关地址，使不同存储形态落到同一键：
    settings 原始形态（``https://host/v1``）与 engine 注入保护形态
    （``https://host/v1/chat/completions#``，fragment 保护，见 engine 交接缝）
    是同一个网关，必须同键，否则能力缓存永不命中（v1.0.39.2 实锤）。
    """
    raw = str(base_url or "").strip()
    if not raw:
        return ""
    raw = raw.rstrip("/")
    if "#" in raw:
        raw = raw.split("#", 1)[0].rstrip("/")
    if raw.endswith("/chat/completions"):
        raw = raw[: -len("/chat/completions")].rstrip("/")
    return raw


def capability_key(provider: Any, base_url: Any, model: Any) -> str:
    """同一网关同一模型的唯一键（缺失字段一律归一为空串）。"""
    return "|".join(
        _normalize_base(v) if i == 1 else str(v or "").strip()
        for i, v in enumerate((provider, base_url, model))
    )


def mark_file_unsupported(provider: Any, base_url: Any, model: Any) -> None:
    """记录该网关不支持 file 块（降级触发后由回调写入）。"""
    if provider is None and base_url is None and model is None:
        return
    _file_unsupported[capability_key(provider, base_url, model)] = True


def file_unsupported(provider: Any, base_url: Any, model: Any) -> bool:
    """该网关是否已知不支持 file 块（打包时命中 → 提前转文本，零 400 往返）。"""
    if not _file_unsupported:
        return False
    return _file_unsupported.get(capability_key(provider, base_url, model), False)


def _decode_file_data(part: dict[str, Any]) -> tuple[str, bytes] | None:
    """从 file 块解出 (mime, 原始字节)。成功返回二元组，失败 None。"""
    file_obj = part.get("file")
    if not isinstance(file_obj, dict):
        return None
    file_data = str(file_obj.get("file_data") or "")
    # 形状：data:<mime>;base64,<b64>
    if not file_data.startswith("data:"):
        return None
    head, sep, b64 = file_data.partition(",")
    if not sep or not b64:
        return None
    mime = head[len("data:"):].split(";", 1)[0].strip().lower() or "application/octet-stream"
    try:
        raw = base64.b64decode(b64, validate=True)
    except Exception:
        return None
    if not raw:
        return None
    return mime, raw


def _extract_with_fitz(raw: bytes) -> str | None:
    """PDF：文字层提取。扫描版（无文字层）返回空字符串 → 上层判为不可提取。"""
    try:
        import fitz  # PyMuPDF
    except Exception:
        return None
    try:
        doc = fitz.open(stream=raw, filetype="pdf")
        try:
            pages = [page.get_text() for page in doc]
        finally:
            doc.close()
        return "\n".join(pages).strip()
    except Exception:
        return None


def _extract_with_docx(raw: bytes) -> str | None:
    try:
        import docx  # python-docx
    except Exception:
        return None
    try:
        document = docx.Document(io.BytesIO(raw))
        return "\n".join(p.text for p in document.paragraphs if p.text).strip()
    except Exception:
        return None


def _extract_with_openpyxl(raw: bytes) -> str | None:
    try:
        import openpyxl
    except Exception:
        return None
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        rows: list[str] = []
        try:
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                    if cells:
                        rows.append("\t".join(cells))
        finally:
            wb.close()
        return "\n".join(rows).strip()
    except Exception:
        return None


def _extract_with_pptx(raw: bytes) -> str | None:
    try:
        from pptx import Presentation  # python-pptx
    except Exception:
        return None
    try:
        prs = Presentation(io.BytesIO(raw))
        chunks: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(getattr(shape, "text_frame", None), "text", "") or ""
                if text.strip():
                    chunks.append(text.strip())
        return "\n".join(chunks).strip()
    except Exception:
        return None


def _extract_with_xlrd(raw: bytes) -> str | None:
    """xls（老二进制格式）。"""
    try:
        import xlrd
    except Exception:
        return None
    try:
        book = xlrd.open_workbook(file_contents=raw)
        rows: list[str] = []
        for sheet in book.sheets():
            for r in range(sheet.nrows):
                cells = []
                for c in range(sheet.ncols):
                    v = sheet.cell_value(r, c)
                    s = str(v).strip() if v is not None else ""
                    if s:
                        cells.append(s)
                if cells:
                    rows.append("\t".join(cells))
        return "\n".join(rows).strip()
    except Exception:
        return None


def extract_file_text(part: dict[str, Any], max_chars: int = EXTRACT_MAX_CHARS) -> str | None:
    """从 file 块提文本。无文字层/格式不支持/解码失败 → None（上层走人话提示）。"""
    decoded = _decode_file_data(part)
    if decoded is None:
        return None
    mime, raw = decoded
    name = str((part.get("file") or {}).get("filename") or "")
    ext = (name.rsplit(".", 1)[-1] if "." in name else "").lower()

    text: str | None = None
    if mime == "application/pdf" or ext == "pdf":
        text = _extract_with_fitz(raw)
    elif mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ) or ext in ("docx", "doc"):
        text = text or _extract_with_docx(raw) or (_extract_with_fitz(raw) if ext == "doc" else None)
    elif mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ) or ext in ("xlsx", "xls"):
        if mime == "application/vnd.ms-excel" or ext == "xls":
            text = _extract_with_xlrd(raw) or _extract_with_openpyxl(raw)
        else:
            text = _extract_with_openpyxl(raw)
    elif mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ) or ext in ("pptx", "ppt"):
        text = _extract_with_pptx(raw)
    elif mime.startswith("text/") or ext in ("txt", "md", "markdown", "csv", "tsv", "json",
                                             "log", "py", "js", "ts", "tsx", "jsx", "html",
                                             "css", "yaml", "yml", "xml", "ini", "cfg", "toml"):
        try:
            text = raw.decode("utf-8", errors="replace").strip()
        except Exception:
            text = None

    if text is None or not text.strip():
        return None
    text = text.strip()
    if max_chars is not None and max_chars > 0 and len(text) > max_chars:
        text = text[:max_chars] + "\n（内容过长，已截断）"
    return text


def replace_file_blocks_with_text(
    messages: list[dict[str, Any]],
    max_chars: int = EXTRACT_MAX_CHARS,
) -> dict[str, Any]:
    """把 messages 里**最后一条 user** 的 file 块逐个换成 text 块（纯函数）。

    返回 {"replaced": n, "failed": [文件名...]}, replaced 即成功换掉的个数。
    图片块（image_url）绝不触碰。多附件混合：失败的文件不阻塞成功的文件。
    """
    result = {"replaced": 0, "failed": []}
    for message in reversed(messages):
        if not isinstance(message, dict) or message.get("role") != "user":
            continue
        content = message.get("content")
        if not isinstance(content, list):
            return result
        new_blocks: list[dict[str, Any]] = []
        for block in content:
            if not isinstance(block, dict) or block.get("type") != "file":
                new_blocks.append(block)
                continue
            text = extract_file_text(block, max_chars=max_chars)
            name = str((block.get("file") or {}).get("filename") or "文件")
            if text is None:
                # [v1.0.39.2-实测修正] 提取失败的文件**换成说明文本块**，绝不保留
                #   file 块原样（残留 file 块 → 降级重发照样 400，白降一次）。
                #   模型读到说明后自然回应「这个文件内容不可见」，用户得到人话。
                result["failed"].append(name)
                new_blocks.append({
                    "type": "text",
                    "text": f"（文件 {name} 已收到，但无法读取内容——可能是扫描版、"
                            f"加密或格式不受支持）",
                })
                continue
            new_blocks.append({
                "type": "text",
                "text": f"【文件 {name} 内容】\n{text}",
            })
            result["replaced"] += 1
        message["content"] = new_blocks
        return result
    return result


def build_format_fallback(
    provider: Any,
    base_url: Any,
    model: Any,
    max_chars: int = EXTRACT_MAX_CHARS,
) -> Any:
    """ProviderClient 用的 on_format_rejected 回调工厂。

    回调契约 (messages) -> list | None：
      · 处理过 file 块（成功转文本 或 失败换成说明块，二者必居其一）→
        返回新 messages（全部是 text 块，ProviderClient 单次重发必然通过），
        并标记能力缓存（同网关后续附件在打包时提前转文本，零 400 往返）；
      · 消息里根本没有 file 块（不应发生）→ 返回 None，照常抛原错。
    """
    def _fallback(messages: list[dict[str, Any]]) -> Any:
        outcome = replace_file_blocks_with_text(messages, max_chars=max_chars)
        if outcome["replaced"] > 0 or outcome["failed"]:
            mark_file_unsupported(provider, base_url, model)
            log.info(
                "附件降级：网关 %s（%s）不支持 file 块 → 已转文本重发，"
                "提取成功 %d 个、说明块 %d 个",
                provider, model, outcome["replaced"], len(outcome["failed"]),
            )
            return messages
        return None

    return _fallback


__all__ = [
    "AttachmentError",
    "build_format_fallback",
    "build_parts",
    "echo_meta",
    "extract_file_text",
    "file_unsupported",
    "inject_into_last_user",
    "load_part",
    "mark_file_unsupported",
    "normalize_path",
    "replace_file_blocks_with_text",
    "sign_path",
    "user_content",
    "verify_path",
]
